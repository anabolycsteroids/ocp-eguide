#!/usr/bin/env python3.12
"""
OCP eGuide — Soutenance Professional Presentation Generator
Creates a polished, collision-free PPTX from scratch.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

# ─── CONSTANTS ──────────────────────────────────────────────────────────
W, H = Inches(13.333), Inches(7.5)  # 16:9 widescreen
EMI_W, EMI_H = 12191695, 6858000

# Colors
OCP_GREEN   = RGBColor(0x00, 0xA0, 0x50)
OCP_GREEN_L = RGBColor(0x00, 0xC8, 0x68)
NAVY        = RGBColor(0x0A, 0x16, 0x28)
DARK_BG     = RGBColor(0x0F, 0x1F, 0x36)
CARD_BG     = RGBColor(0x14, 0x2A, 0x4A)
CARD_BG2    = RGBColor(0x1A, 0x34, 0x58)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY    = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY   = RGBColor(0x66, 0x66, 0x66)
SUBTLE_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
AMBER       = RGBColor(0xFF, 0xB3, 0x00)
RED_ACCENT  = RGBColor(0xE8, 0x3E, 0x3E)

SLIDE_BG    = RGBColor(0x0B, 0x1A, 0x2E)
ACCENT_LINE = OCP_GREEN

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

MARGIN_L = 0.8
MARGIN_R = 0.8
CONTENT_W = SLIDE_W_IN - MARGIN_L - MARGIN_R

def inches(v):
    return Inches(v)

def emu(v):
    return int(v * 914400)

def add_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SLIDE_BG
    return slide

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_w or 1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Set corner radius smaller
    sp = shape._element
    prstGeom = sp.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = sp.makeelement(qn('a:avLst'), {})
            prstGeom.append(avLst)
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd = sp.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': 'val 8000'})
        avLst.append(gd)
    return shape

def add_text_box(slide, x, y, w, h, text, font_size=14, bold=False,
                 color=WHITE, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name="Calibri", line_spacing=None):
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        txBox.anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox

def add_multi_text(slide, x, y, w, h, paragraphs, default_size=14,
                   default_color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT,
                   default_bold=False, line_spacing=None, font_name="Calibri"):
    """paragraphs: list of (text, size, color, bold) or just text strings"""
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(paragraphs):
        if isinstance(item, str):
            text, sz, clr, bld = item, default_size, default_color, default_bold
        else:
            text = item[0]
            sz = item[1] if len(item) > 1 else default_size
            clr = item[2] if len(item) > 2 else default_color
            bld = item[3] if len(item) > 3 else default_bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.alignment = alignment
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(sz)
        run.font.bold = bld
        run.font.color.rgb = clr
        run.font.name = font_name
        # Space after paragraph
        p.space_after = Pt(4)
    return txBox

def add_footer(slide, num, total):
    """Minimal footer: 'OCP eGuide  •  2026' + page number"""
    add_rect(slide, 0, 7.08, SLIDE_W_IN, 0.42, RGBColor(0x08, 0x12, 0x22))
    add_text_box(slide, 0.5, 7.14, 6, 0.3, "OCP eGuide  •  2026",
                 font_size=9, color=MID_GRAY, font_name="Calibri")
    add_text_box(slide, 11.5, 7.14, 1.5, 0.3, f"{num} / {total}",
                 font_size=9, color=MID_GRAY, alignment=PP_ALIGN.RIGHT, font_name="Calibri")

def add_green_accent_line(slide, x, y, w):
    add_rect(slide, x, y, w, 0.04, OCP_GREEN)

def add_section_number(slide, num, x=0.8, y=0.35):
    """Big section number in green"""
    add_text_box(slide, x, y, 1.2, 0.7, num,
                 font_size=32, bold=True, color=OCP_GREEN, font_name="Calibri Light")

def add_slide_title(slide, title, subtitle=None, x=0.8, y=1.1, subtitle_y=None):
    """Add slide title and optional subtitle"""
    add_text_box(slide, x, y, 11, 0.7, title,
                 font_size=28, bold=True, color=WHITE, font_name="Calibri Light")
    add_green_accent_line(slide, x, y + 0.72, 2.5)
    if subtitle:
        sy = subtitle_y or (y + 0.9)
        add_text_box(slide, x, sy, 11, 0.5, subtitle,
                     font_size=14, color=MID_GRAY, font_name="Calibri")

def add_stat_card(slide, x, y, w, h, number, label, accent=OCP_GREEN):
    """Visual statistic card"""
    card = add_rounded_rect(slide, x, y, w, h, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
    # Top accent bar
    add_rect(slide, x + 0.05, y + 0.05, w - 0.1, 0.04, accent)
    # Number
    add_text_box(slide, x, y + 0.25, w, 0.6, number,
                 font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
                 font_name="Calibri Light")
    # Label
    add_text_box(slide, x, y + 0.85, w, 0.5, label,
                 font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name="Calibri")
    return card

def add_feature_card(slide, x, y, w, h, icon, title, description, accent=OCP_GREEN):
    """Feature card with icon, title, description"""
    card = add_rounded_rect(slide, x, y, w, h, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, inches(x + 0.2), inches(y + 0.2),
                                      inches(0.45), inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    # Icon text
    add_text_box(slide, x + 0.2, y + 0.22, 0.45, 0.45, icon,
                 font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + 0.8, y + 0.25, w - 1.0, 0.35, title,
                 font_size=13, bold=True, color=WHITE, font_name="Calibri")
    # Description
    add_text_box(slide, x + 0.2, y + 0.75, w - 0.4, h - 0.9, description,
                 font_size=10.5, color=LIGHT_GRAY, font_name="Calibri")
    return card

def add_flow_step(slide, x, y, num, title, desc, w=2.3, accent=OCP_GREEN):
    """Pipeline step with number circle"""
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, inches(x + (w-0.5)/2), inches(y),
                                      inches(0.5), inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text_box(slide, x + (w-0.5)/2, y + 0.05, 0.5, 0.4, str(num),
                 font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x, y + 0.65, w, 0.35, title,
                 font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
                 font_name="Calibri")
    # Description
    add_text_box(slide, x, y + 1.0, w, 1.2, desc,
                 font_size=9.5, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name="Calibri")

def add_arrow_right(slide, x, y):
    add_text_box(slide, x, y, 0.4, 0.4, "→",
                 font_size=22, bold=True, color=OCP_GREEN, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# MAIN GENERATION
# ════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = emu(SLIDE_W_IN)
prs.slide_height = emu(SLIDE_H_IN)

TOTAL_SLIDES = 20

# ────────────────────────────────────────────────────────────────────────
# SLIDE 1: TITLE
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
# Large dark background block (top 45%)
add_rect(s, 0, 0, SLIDE_W_IN, 3.4, NAVY)
# Green accent strip
add_rect(s, 0, 3.4, SLIDE_W_IN, 0.06, OCP_GREEN)

# Project name - huge
add_text_box(s, 1.0, 0.8, 11, 0.9, "OCP eGuide",
             font_size=44, bold=True, color=WHITE, font_name="Calibri Light")

# Subtitle
add_text_box(s, 1.0, 1.7, 11, 0.6, "Système Numérique Interactif d'Orientation du Campus",
             font_size=20, color=OCP_GREEN, font_name="Calibri")

# Tagline
add_text_box(s, 1.0, 2.4, 11, 0.5, "Carte Interactive  •  Routage A*  •  Dashboards  •  Trilingue",
             font_size=13, color=MID_GRAY, font_name="Calibri")

# Info cards in bottom section
info_items = [
    ("Imrane Belkoufa", "Étudiant — Génie Informatique"),
    ("EST Sidi Bennour", "École Supérieure de Technologie"),
    ("OCP S.A. — Jorf Lasfar", "Site industriel, El Jadida"),
    ("M. Karim Alaoui", "Encadrant de stage"),
]
card_w = 2.6
gap = 0.27
start_x = 1.0
for i, (title, sub) in enumerate(info_items):
    cx = start_x + i * (card_w + gap)
    add_rounded_rect(s, cx, 3.9, card_w, 1.4, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
    add_text_box(s, cx + 0.2, 4.05, card_w - 0.4, 0.4, title,
                 font_size=12, bold=True, color=WHITE)
    add_text_box(s, cx + 0.2, 4.45, card_w - 0.4, 0.7, sub,
                 font_size=10, color=MID_GRAY)

# Year
add_text_box(s, 1.0, 5.8, 11, 0.5, "Année universitaire  2025 – 2026",
             font_size=14, color=LIGHT_GRAY, font_name="Calibri")

# Footer
add_footer(s, 1, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 2: PLAN DE LA PRÉSENTATION
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "01")
add_slide_title(s, "Plan de la Présentation")

plan_items = [
    ("Contexte & Problématique", "Complexe Jorf Lasfar"),
    ("Solution OCP eGuide", "Fonctionnalités clés"),
    ("Architecture & Stack", "Frontend → Backend → BDD"),
    ("Carte & Calibration", "Canvas HTML5 + GPS"),
    ("Pathfinding A*", "Routage piéton optimal"),
    ("Backend & API", "Express + Sécurité"),
    ("Dashboards", "19 profils utilisateurs"),
    ("Développement", "Processus & outils"),
    ("Tests & Validation", "Qualité logicielle"),
    ("Déploiement", "Vercel + Render + Supabase"),
    ("Résultats", "Métriques clés"),
    ("Défis Techniques", "Problèmes résolus"),
    ("Perspectives", "Évolutions futures"),
    ("Conclusion", "Synthèse"),
]

cols = 2
rows_per_col = 7
card_w = 5.4
card_h = 0.52
x_starts = [1.0, 7.0]
y_start = 2.1
y_gap = 0.62

for i, (title, sub) in enumerate(plan_items):
    col = i // rows_per_col
    row = i % rows_per_col
    cx = x_starts[col]
    cy = y_start + row * y_gap
    
    # Number badge
    badge = add_rounded_rect(s, cx, cy, 0.38, 0.38, OCP_GREEN)
    add_text_box(s, cx, cy + 0.03, 0.38, 0.32, f"{i+1:02d}",
                 font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Title
    add_text_box(s, cx + 0.5, cy - 0.02, 4.0, 0.3, title,
                 font_size=12, bold=True, color=WHITE)
    # Subtitle
    add_text_box(s, cx + 0.5, cy + 0.24, 4.0, 0.25, sub,
                 font_size=9, color=MID_GRAY)

add_footer(s, 2, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 3: CONTEXTE — OCP S.A.
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "02")
add_slide_title(s, "Contexte : OCP S.A.")

# Left column: OCP S.A.
add_rounded_rect(s, 0.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.1, 2.3, 5.0, 0.4, "OCP S.A.",
             font_size=18, bold=True, color=OCP_GREEN, font_name="Calibri")
add_green_accent_line(s, 1.1, 2.75, 1.5)

stats_left = [
    "Leader mondial dans l'exploitation et la\ntransformation des phosphates",
    "Groupe industriel marocain,\nprésence dans 15+ pays",
    "Fournit 40 % de l'engrais mondial",
    "Chiffre d'affaires : 13,3 milliards MAD",
    "Plus de 21 000 collaborateurs",
]
for i, txt in enumerate(stats_left):
    add_text_box(s, 1.3, 3.0 + i * 0.72, 4.8, 0.7, f"▸  {txt}",
                 font_size=11, color=LIGHT_GRAY, line_spacing=15)

# Right column: Complexe Jorf Lasfar
add_rounded_rect(s, 6.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 7.1, 2.3, 5.0, 0.4, "Complexe Jorf Lasfar",
             font_size=18, bold=True, color=OCP_GREEN, font_name="Calibri")
add_green_accent_line(s, 7.1, 2.75, 2.0)

stats_right = [
    "Situé près d'El Jadida,\ncôte atlantique marocaine",
    "1 130 hectares de\nsuperficie totale",
    "Plus de 90 localisations\nà orienter",
    "19 catégories de\nprofils utilisateurs",
    "Environnement industriel\ncomplexe et étendu",
]
for i, txt in enumerate(stats_right):
    add_text_box(s, 7.3, 3.0 + i * 0.72, 4.8, 0.7, f"▸  {txt}",
                 font_size=11, color=LIGHT_GRAY, line_spacing=15)

add_footer(s, 3, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 4: PROBLÉMATIQUE
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "02")
add_slide_title(s, "Problématique : L'Orientation sur Site")

# Stat cards row
stat_data = [
    ("90+", "localisations\nà orienter"),
    ("19", "profils\nd'utilisateurs"),
    ("1 130 ha", "de superficie\nindustrielle"),
    ("0", "solution\nnumérique existante"),
]
for i, (num, lbl) in enumerate(stat_data):
    add_stat_card(s, 0.8 + i * 3.1, 2.1, 2.8, 1.5, num, lbl)

# Problem description card
add_rounded_rect(s, 0.8, 4.0, 11.7, 2.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.1, 4.2, 11, 0.4, "Le Problème",
             font_size=16, bold=True, color=RED_ACCENT, font_name="Calibri")
add_green_accent_line(s, 1.1, 4.65, 1.2)

problems = [
    "Plans imprimés obsolètes — ne reflètent pas les modifications récentes du site",
    "GPS grand public imprécis en milieu industriel (obstructions, multi-trajet)",
    "Absence de solution adaptée aux spécificités du campus OCP Jorf Lasfar",
    "Difficulté d'orientation pour les nouveaux employés, stagiaires et visiteurs",
]
for i, prob in enumerate(problems):
    add_text_box(s, 1.3, 4.85 + i * 0.42, 10.5, 0.4, f"▸  {prob}",
                 font_size=12, color=LIGHT_GRAY)

add_footer(s, 4, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 5: SOLUTION OCP eGuide
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "03")
add_slide_title(s, "Solution : OCP eGuide")

# Intro text
add_text_box(s, 0.8, 2.0, 11.5, 0.5,
             "Application web complète : carte interactive, recherche floue, routage A*, dashboards, mode kiosk, trilingue.",
             font_size=13, color=LIGHT_GRAY)

# Feature cards - 3x2 grid
features = [
    ("🗺", "Carte Interactive", "Image satellite du campus avec zoom,\npan et marqueurs géolocalisés"),
    ("🔍", "Recherche Floue", "Recherche par nom, alias ou catégorie\navec scoring multi-critères"),
    ("🧭", "Routage A*", "Itinéraire piéton optimal avec\ntemps de marche et instructions"),
    ("📊", "19 Dashboards", "Interfaces personnalisées par profil\nutilisateur avec données pertinentes"),
    ("🌐", "Trilingue", "Français, anglais, arabe avec\nadaptation RTL — 430+ clés i18n"),
    ("🔒", "Sécurisé", "JWT, Argon2, rate limiting,\nCORS, Helmet.js, validation Zod"),
]
card_w = 3.7
card_h = 1.8
gap_x = 0.35
gap_y = 0.3
x_start = 0.8
y_start = 2.5

for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    cx = x_start + col * (card_w + gap_x)
    cy = y_start + row * (card_h + gap_y)
    
    card = add_rounded_rect(s, cx, cy, card_w, card_h, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
    
    # Icon background
    icon_bg = add_rounded_rect(s, cx + 0.2, cy + 0.15, 0.45, 0.45, OCP_GREEN)
    add_text_box(s, cx + 0.2, cy + 0.17, 0.45, 0.4, icon,
                 font_size=16, alignment=PP_ALIGN.CENTER, color=WHITE)
    
    # Title
    add_text_box(s, cx + 0.8, cy + 0.18, card_w - 1.0, 0.3, title,
                 font_size=13, bold=True, color=WHITE)
    
    # Description
    add_text_box(s, cx + 0.2, cy + 0.7, card_w - 0.4, card_h - 0.85, desc,
                 font_size=10.5, color=LIGHT_GRAY)

add_footer(s, 5, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 6: ARCHITECTURE DU SYSTÈME
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "03")
add_slide_title(s, "Architecture du Système")

# Three-tier architecture visual
layers = [
    ("PRÉSENTATION", "Next.js 16  •  React 19  •  TypeScript  •  Tailwind CSS  •  Canvas HTML5",
     OCP_GREEN, 2.2),
    ("MÉTIER", "Express.js 4  •  Socket.IO  •  JWT  •  Zod  •  Swagger  •  Rate Limiting",
     RGBColor(0x3A, 0x8A, 0xD6), 3.5),
    ("DONNÉES", "PostgreSQL 16  •  Prisma ORM  •  Argon2  •  Supabase",
     RGBColor(0xE8, 0x8A, 0x2E), 4.8),
]

for label, techs, color, y in layers:
    # Layer card
    add_rounded_rect(s, 1.5, y, 10.3, 1.0, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
    # Color accent bar on left
    add_rect(s, 1.5, y, 0.08, 1.0, color)
    # Label
    add_text_box(s, 1.8, y + 0.12, 3.0, 0.35, label,
                 font_size=14, bold=True, color=color)
    # Technologies
    add_text_box(s, 1.8, y + 0.52, 9.5, 0.35, techs,
                 font_size=12, color=LIGHT_GRAY)

# Arrows between layers
add_text_box(s, 6.2, 3.2, 1, 0.3, "▼",
             font_size=20, bold=True, color=OCP_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(s, 6.2, 4.5, 1, 0.3, "▼",
             font_size=20, bold=True, color=OCP_GREEN, alignment=PP_ALIGN.CENTER)

# Deployment line
add_rounded_rect(s, 0.8, 6.2, 11.7, 0.6, RGBColor(0x0D, 0x1E, 0x35), border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.0, 6.28, 11, 0.4,
             "Déploiement  :  Vercel (frontend)  •  Render (backend)  •  Supabase (PostgreSQL)",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_footer(s, 6, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 7: STACK TECHNOLOGIQUE
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "03")
add_slide_title(s, "Stack Technologique")

# Table-like grid with cards
tech_data = [
    ("Catégorie", "Technologie", "Version / Détails"),
    ("Frontend", "Next.js + React", "v16 / v19 — SSR, App Router"),
    ("Langage", "TypeScript", "Type safety end-to-end"),
    ("Styling", "Tailwind CSS", "Utility-first, responsive"),
    ("Rendu carte", "Canvas HTML5", "DPR-aware, zoom 0.3–5.0x"),
    ("Backend", "Express.js", "v4 — REST API, middleware"),
    ("Temps réel", "Socket.IO", "Notifications, présence"),
    ("Auth", "JWT + Argon2", "Access (15 min) + Refresh (7 j)"),
    ("Validation", "Zod", "Schémas de validation stricte"),
    ("ORM", "Prisma", "16 modèles, type-safe queries"),
    ("Base de données", "PostgreSQL", "v16 — Supabase hosting"),
    ("Docs API", "Swagger", "Auto-générée depuis les routes"),
    ("i18n", "Custom Provider", "3 langues, 430+ clés, RTL"),
]

row_h = 0.37
x_start = 0.8
y_start = 2.0
col_widths = [2.0, 2.5, 7.2]
col_x = [x_start, x_start + col_widths[0], x_start + col_widths[0] + col_widths[1]]

for row_idx, row_data in enumerate(tech_data):
    y = y_start + row_idx * row_h
    is_header = row_idx == 0
    bg_color = RGBColor(0x0D, 0x1E, 0x35) if is_header else (CARD_BG if row_idx % 2 == 1 else RGBColor(0x10, 0x22, 0x3C))
    
    for col_idx, text in enumerate(row_data):
        add_rect(s, col_x[col_idx], y, col_widths[col_idx], row_h, bg_color,
                 border_color=RGBColor(0x1E, 0x3A, 0x5C))
        font_color = OCP_GREEN if is_header else (WHITE if col_idx <= 1 else LIGHT_GRAY)
        font_bold = is_header or col_idx == 0
        add_text_box(s, col_x[col_idx] + 0.15, y + 0.04, col_widths[col_idx] - 0.3, row_h - 0.08,
                     text, font_size=10 if not is_header else 10.5, bold=font_bold, color=font_color)

add_footer(s, 7, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 8: CARTE INTERACTIVE DU CAMPUS
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "04")
add_slide_title(s, "Carte Interactive du Campus")

# Left: why canvas
add_rounded_rect(s, 0.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.1, 2.3, 5.0, 0.4, "Pourquoi Canvas HTML5 ?",
             font_size=16, bold=True, color=OCP_GREEN, font_name="Calibri")
add_green_accent_line(s, 1.1, 2.75, 1.8)

canvas_reasons = [
    "Image satellite PNG comme surface\nde référence unique",
    "Positionnement au pixel-près des\nmarqueurs à tout niveau de zoom",
    "Pas de dépendance à Leaflet/Mapbox\n(coût, licence, restriction)",
    "Rendu haute performance pour\n95 localisations simultanées",
    "DPR-aware : qualité optimale sur\nmobile, tablette et desktop",
]
for i, txt in enumerate(canvas_reasons):
    add_text_box(s, 1.3, 3.0 + i * 0.72, 4.8, 0.7, f"▸  {txt}",
                 font_size=11, color=LIGHT_GRAY, line_spacing=15)

# Right: image placeholder
add_rounded_rect(s, 6.8, 2.1, 5.6, 4.6, RGBColor(0x08, 0x12, 0x22), border_color=RGBColor(0x1E, 0x3A, 0x5C))
# Try to add campus map image
campus_map_path = "freebuff-frontend/public/assets/map/campus-map.png"
if os.path.exists(campus_map_path):
    try:
        s.shapes.add_picture(campus_map_path, inches(7.0), inches(2.3), inches(5.2), inches(3.2))
        add_text_box(s, 7.0, 5.6, 5.2, 0.4, "Vue principale de la carte interactive",
                     font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    except:
        add_text_box(s, 7.5, 4.0, 4.5, 0.5, "[Image de la carte campus]",
                     font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
else:
    add_text_box(s, 7.5, 4.0, 4.5, 0.5, "[Image de la carte campus]",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer(s, 8, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 9: CALIBRATION GPS → PIXELS
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "04")
add_slide_title(s, "Calibration GPS → Pixels")

# Pipeline flow
steps = [
    ("1", "Capture", "8 bâtiments identifiés\nsur satellite +\nGoogle Earth"),
    ("2", "Calcul", "Transformation affine\n6 coefficients\n(a,b,c,d,e,f)"),
    ("3", "Application", "Conversion des 95\nlocalisations GPS\n→ pixels"),
    ("4", "Validation", "RMS = 14,4 px\n95/95 in-bounds\n100 % routable"),
]

step_w = 2.4
arrow_w = 0.6
total_w = len(steps) * step_w + (len(steps) - 1) * arrow_w
start_x = (SLIDE_W_IN - total_w) / 2

for i, (num, title, desc) in enumerate(steps):
    sx = start_x + i * (step_w + arrow_w)
    add_flow_step(s, sx, 2.2, num, title, desc, w=step_w)
    if i < len(steps) - 1:
        add_arrow_right(s, sx + step_w + 0.1, 2.35)

# Result callout
add_rounded_rect(s, 2.5, 5.2, 8.3, 1.2, CARD_BG, border_color=OCP_GREEN)
add_rect(s, 2.5, 5.2, 8.3, 0.05, OCP_GREEN)
add_text_box(s, 2.8, 5.4, 7.7, 0.8,
             "Calibration validée : 14,4 px (≈ 100 m)  —  95/95 in-bounds  —  100 % routable",
             font_size=16, bold=True, color=OCP_GREEN, alignment=PP_ALIGN.CENTER)

add_footer(s, 9, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 10: PATHFINDING A*
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "05")
add_slide_title(s, "Pathfinding A* et Routage Piéton")

# Left: Algorithm explanation
add_rounded_rect(s, 0.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.1, 2.3, 5.0, 0.4, "Algorithme A*",
             font_size=16, bold=True, color=OCP_GREEN, font_name="Calibri")
add_green_accent_line(s, 1.1, 2.75, 1.2)

algo_lines = [
    ("f(n) = g(n) + h(n)", 14, OCP_GREEN, True),
    ("", 6, LIGHT_GRAY, False),
    ("g(n) : coût réel du début au nœud n", 12, LIGHT_GRAY, False),
    ("h(n) : heuristique euclidienne vers la destination", 12, LIGHT_GRAY, False),
    ("", 6, LIGHT_GRAY, False),
    ("Optimisations :", 13, WHITE, True),
    ("▸  File de priorité (tas binaire)", 12, LIGHT_GRAY, False),
    ("▸  Adjacence paresseuse pour la mémoire", 12, LIGHT_GRAY, False),
    ("▸  Conversion pixels → mètres (1,4 m/s)", 12, LIGHT_GRAY, False),
    ("", 6, LIGHT_GRAY, False),
    ("Pipeline :", 13, WHITE, True),
    ("Lieu → place-node-connections → nœud A*", 11, LIGHT_GRAY, False),
    ("→ chemin optimal → instructions", 11, LIGHT_GRAY, False),
]
add_multi_text(s, 1.3, 3.0, 4.8, 3.5, algo_lines, font_name="Calibri")

# Right: Statistics
add_rounded_rect(s, 6.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 7.1, 2.3, 5.0, 0.4, "Données Cartographiques",
             font_size=16, bold=True, color=OCP_GREEN, font_name="Calibri")
add_green_accent_line(s, 7.1, 2.75, 2.0)

route_stats = [
    ("3 423", "nœuds de navigation"),
    ("13 346", "arêtes (chemins piétons)"),
    ("95", "localisations connectées"),
    ("100 %", "paires routables"),
    ("< 3 s", "temps de calcul max"),
]
for i, (num, lbl) in enumerate(route_stats):
    row_y = 3.1 + i * 0.72
    add_text_box(s, 7.3, row_y, 1.8, 0.4, num,
                 font_size=22, bold=True, color=WHITE, font_name="Calibri Light")
    add_text_box(s, 9.2, row_y + 0.05, 3.0, 0.35, lbl,
                 font_size=12, color=LIGHT_GRAY)

add_footer(s, 10, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 11: BACKEND, API & BASE DE DONNÉES
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "06")
add_slide_title(s, "Backend, API et Base de Données")

# API groups table
api_data = [
    ("Groupe de routes", "Fonctionnalité", "Sécurité"),
    ("Auth", "Login, refresh, logout, register", "JWT + Argon2"),
    ("Users", "CRUD utilisateurs, profils", "Auth + Rôle"),
    ("Guests", "API publique, QR codes", "Token invite"),
    ("Locations", "90+ localisations, recherche", "Public (GET)"),
    ("Roles", "19 profils, permissions", "Admin"),
    ("Profiles", "Profils par catégorie", "Auth"),
    ("Map Data", "nœuds, arêtes, connexions", "Public (GET)"),
    ("Calibration", "GPS → pixels, validation", "Admin"),
    ("Dashboard", "Stats, widgets temps réel", "Auth + Rôle"),
    ("Notifications", "Push temps réel (Socket.IO)", "Auth"),
    ("Export", "PDF, CSV des rapports", "Auth + Rôle"),
    ("Config", "Paramètres système", "Admin"),
]

row_h = 0.35
x_start = 0.8
y_start = 2.0
col_widths = [2.0, 5.2, 4.3]
col_x = [x_start, x_start + col_widths[0], x_start + col_widths[0] + col_widths[1]]

for row_idx, row_data in enumerate(api_data):
    y = y_start + row_idx * row_h
    is_header = row_idx == 0
    bg_color = RGBColor(0x0D, 0x1E, 0x35) if is_header else (CARD_BG if row_idx % 2 == 1 else RGBColor(0x10, 0x22, 0x3C))
    
    for col_idx, text in enumerate(row_data):
        add_rect(s, col_x[col_idx], y, col_widths[col_idx], row_h, bg_color,
                 border_color=RGBColor(0x1E, 0x3A, 0x5C))
        font_color = OCP_GREEN if is_header else (WHITE if col_idx <= 1 else LIGHT_GRAY)
        add_text_box(s, col_x[col_idx] + 0.12, y + 0.03, col_widths[col_idx] - 0.24, row_h - 0.06,
                     text, font_size=9.5, bold=is_header, color=font_color)

# Bottom summary
add_text_box(s, 0.8, 6.8, 11, 0.3,
             "12 groupes de routes  •  Documentation Swagger auto-générée  •  16 modèles Prisma",
             font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer(s, 11, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 12: SÉCURITÉ
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "06")
add_slide_title(s, "Sécurité")

security_items = [
    ("🔑", "Authentification JWT", "Double token : access (15 min) + refresh (7 jours)"),
    ("🛡", "Hachage Argon2", "Algorithme de hachage sécurisé pour les mots de passe"),
    ("👤", "Autorisation par rôle", "Middleware authorize : ADMIN, EMPLOYEE, INTERN"),
    ("⏱", "Rate Limiting", "100 requêtes / 15 min par IP"),
    ("🔒", "Helmet.js", "En-têtes de sécurité HTTP renforcés"),
    ("🌐", "CORS", "Liste blanche d'origines autorisées"),
    ("✅", "Validation Zod", "Validation stricte de toutes les entrées API"),
    ("🧹", "Sanitisation", "Suppression du passwordHash dans les réponses"),
]

cols = 2
rows = 4
card_w = 5.6
card_h = 0.95
gap_x = 0.35
gap_y = 0.2
x_start = 0.8
y_start = 2.1

for i, (icon, title, desc) in enumerate(security_items):
    col = i % cols
    row = i // cols
    cx = x_start + col * (card_w + gap_x)
    cy = y_start + row * (card_h + gap_y)
    
    add_rounded_rect(s, cx, cy, card_w, card_h, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
    
    # Icon
    icon_bg = add_rounded_rect(s, cx + 0.15, cy + 0.2, 0.45, 0.45, RGBColor(0x1A, 0x3A, 0x5C))
    add_text_box(s, cx + 0.15, cy + 0.22, 0.45, 0.4, icon,
                 font_size=16, alignment=PP_ALIGN.CENTER, color=WHITE)
    
    # Title
    add_text_box(s, cx + 0.75, cy + 0.15, card_w - 1.0, 0.35, title,
                 font_size=12, bold=True, color=WHITE)
    # Description
    add_text_box(s, cx + 0.75, cy + 0.5, card_w - 1.0, 0.35, desc,
                 font_size=10, color=LIGHT_GRAY)

add_footer(s, 12, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 13: INTERFACE & DASHBOARDS
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "07")
add_slide_title(s, "Interface Utilisateur et Dashboards")

# Left: identity
add_rounded_rect(s, 0.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.1, 2.3, 5.0, 0.4, "Identité Visuelle",
             font_size=16, bold=True, color=OCP_GREEN, font_name="Calibri")
add_green_accent_line(s, 1.1, 2.75, 1.5)

ui_features = [
    "Couleurs : Vert OCP (#00A050) + Marine (#0A1628)",
    "Mode kiosk : défilement automatique des profils",
    "Page d'accueil : sélection de profil avec icônes",
    "Responsive design : mobile, tablette, desktop",
    "i18n : 3 langues (EN/FR/AR), 430+ clés, RTL",
]
for i, txt in enumerate(ui_features):
    add_text_box(s, 1.3, 3.0 + i * 0.65, 4.8, 0.6, f"▸  {txt}",
                 font_size=11, color=LIGHT_GRAY, line_spacing=15)

# Right: Dashboard categories
add_rounded_rect(s, 6.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 7.1, 2.3, 5.0, 0.4, "19 Dashboards Personnalisés",
             font_size=16, bold=True, color=OCP_GREEN, font_name="Calibri")
add_green_accent_line(s, 7.1, 2.75, 2.5)

dash_categories = [
    ("Employés", "Management, Réception, RH, IT, Sécurité"),
    ("Stagiaires", "Mécanique, Chimie, Électrique, Civil..."),
    ("Visiteurs", "Client, Livreur, Partenaire, Fournisseur..."),
]
for i, (cat, details) in enumerate(dash_categories):
    row_y = 3.0 + i * 0.95
    add_rounded_rect(s, 7.1, row_y, 5.0, 0.8, RGBColor(0x0D, 0x1E, 0x35))
    add_text_box(s, 7.3, row_y + 0.05, 4.6, 0.35, cat,
                 font_size=13, bold=True, color=WHITE)
    add_text_box(s, 7.3, row_y + 0.4, 4.6, 0.35, details,
                 font_size=10, color=LIGHT_GRAY)

# Widget note
add_text_box(s, 7.1, 5.9, 5.0, 0.4,
             "▸  Widget temps réel via Socket.IO",
             font_size=11, color=LIGHT_GRAY)

add_footer(s, 13, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 14: PROCESSUS DE DÉVELOPPEMENT
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "08")
add_slide_title(s, "Processus de Développement")

phases = [
    ("1", "Analyse", "Étude du site,\nbesoins utilisateurs,\n90+ localisations"),
    ("2", "Conception", "Architecture 3-tiers,\nschéma de base,\nwireframes"),
    ("3", "Développement", "Frontend Next.js,\nBackend Express,\nCanvas map"),
    ("4", "Calibration", "GPS → pixels,\n8 points de contrôle,\ntransformation affine"),
    ("5", "Intégration", "A*, Socket.IO,\ni18n, dashboards,\nsécurité"),
    ("6", "Validation", "Tests TypeScript,\n Routage 100 %,\n déploiement"),
]

step_w = 1.8
gap = 0.25
total = len(phases) * step_w + (len(phases) - 1) * gap
start_x = (SLIDE_W_IN - total) / 2

for i, (num, title, desc) in enumerate(phases):
    sx = start_x + i * (step_w + gap)
    
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, inches(sx + (step_w - 0.55) / 2), inches(2.2),
                                  inches(0.55), inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = OCP_GREEN
    circle.line.fill.background()
    add_text_box(s, sx + (step_w - 0.55) / 2, 2.24, 0.55, 0.48, num,
                 font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Connector line
    if i < len(phases) - 1:
        add_rect(s, sx + step_w, 2.45, gap, 0.04, RGBColor(0x1E, 0x3A, 0x5C))
    
    # Title
    add_text_box(s, sx, 2.9, step_w, 0.35, title,
                 font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Description
    add_text_box(s, sx, 3.3, step_w, 1.2, desc,
                 font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Tools section
add_rounded_rect(s, 0.8, 5.0, 11.7, 1.5, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.1, 5.15, 11, 0.35, "Outils utilisés",
             font_size=14, bold=True, color=WHITE, font_name="Calibri")

tools = [
    ("VS Code", "Éditeur principal"),
    ("Git / GitHub", "Version control"),
    ("Google Earth", "Calibration GPS"),
    ("Postman", "Tests API"),
    ("Prisma Studio", "Gestion BDD"),
    ("Vercel / Render", "Déploiement"),
]
tool_w = 1.8
tool_gap = 0.15
for i, (tool, usage) in enumerate(tools):
    tx = 1.1 + i * (tool_w + tool_gap)
    add_text_box(s, tx, 5.6, tool_w, 0.3, tool,
                 font_size=11, bold=True, color=OCP_GREEN)
    add_text_box(s, tx, 5.95, tool_w, 0.3, usage,
                 font_size=9, color=MID_GRAY)

add_footer(s, 14, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 15: TESTS & VALIDATION
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "09")
add_slide_title(s, "Tests et Validation")

test_cards = [
    ("✓", "TypeScript", "0 erreur de compilation\nVérification de types\nstricte end-to-end", OCP_GREEN),
    ("🗺", "Carte", "95/95 localisations\nin-bounds\nLabels lisibles\nà tout zoom", RGBColor(0x3A, 0x8A, 0xD6)),
    ("🔍", "Recherche", "90+ lieux trouvés\nFonctionnement des aliases\nScoring multi-critères", RGBColor(0xE8, 0x8A, 0x2E)),
    ("🧭", "Routage A*", "100 % paires routables\nInstructions générées\nTemps < 3 s", OCP_GREEN),
    ("🔒", "Sécurité", "JWT fonctionnel\n76 comptes testés\n19 profils validés", RGBColor(0xE8, 0x3E, 0x3E)),
    ("🌐", "i18n", "3 langues (EN/FR/AR)\n430+ clés traduites\nRTL arabe validé", RGBColor(0x9B, 0x59, 0xB6)),
]

card_w = 3.7
card_h = 2.5
gap_x = 0.35
gap_y = 0.3
x_start = 0.8
y_start = 2.1

for i, (icon, title, desc, accent) in enumerate(test_cards):
    col = i % 3
    row = i // 3
    cx = x_start + col * (card_w + gap_x)
    cy = y_start + row * (card_h + gap_y)
    
    add_rounded_rect(s, cx, cy, card_w, card_h, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
    # Top accent
    add_rect(s, cx + 0.05, cy + 0.05, card_w - 0.1, 0.04, accent)
    
    # Icon
    icon_bg = add_rounded_rect(s, cx + 0.2, cy + 0.25, 0.5, 0.5, accent)
    add_text_box(s, cx + 0.2, cy + 0.27, 0.5, 0.45, icon,
                 font_size=18, alignment=PP_ALIGN.CENTER, color=WHITE)
    
    # Title
    add_text_box(s, cx + 0.85, cy + 0.3, card_w - 1.1, 0.35, title,
                 font_size=14, bold=True, color=WHITE)
    
    # Description
    add_text_box(s, cx + 0.2, cy + 0.9, card_w - 0.4, card_h - 1.1, desc,
                 font_size=11, color=LIGHT_GRAY)

add_footer(s, 15, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 16: DÉPLOIEMENT
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "10")
add_slide_title(s, "Déploiement et Production")

# Deployment flow
deploy_steps = [
    ("Vercel", "Frontend\nNext.js SSR\nDéploiement automatique\nSSL intégré", OCP_GREEN),
    ("Render", "Backend\nExpress.js\nAuto-scaling\nLogs & monitoring", RGBColor(0x3A, 0x8A, 0xD6)),
    ("Supabase", "Base de données\nPostgreSQL 16\nBackups automatiques\nPanneau de contrôle", RGBColor(0xE8, 0x8A, 0x2E)),
]

step_w = 3.4
gap = 0.5
total = len(deploy_steps) * step_w + (len(deploy_steps) - 1) * gap
start_x = (SLIDE_W_IN - total) / 2

for i, (title, desc, accent) in enumerate(deploy_steps):
    sx = start_x + i * (step_w + gap)
    
    add_rounded_rect(s, sx, 2.2, step_w, 3.5, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
    # Top accent
    add_rect(s, sx + 0.05, 2.25, step_w - 0.1, 0.05, accent)
    
    # Title
    add_text_box(s, sx, 2.5, step_w, 0.5, title,
                 font_size=20, bold=True, color=accent, alignment=PP_ALIGN.CENTER,
                 font_name="Calibri Light")
    
    add_green_accent_line(s, sx + 0.5, 3.05, step_w - 1.0)
    
    # Description
    add_text_box(s, sx + 0.3, 3.25, step_w - 0.6, 2.2, desc,
                 font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    # Arrows
    if i < len(deploy_steps) - 1:
        add_arrow_right(s, sx + step_w + 0.05, 3.6)

# Bottom note
add_rounded_rect(s, 1.5, 6.0, 10.3, 0.65, RGBColor(0x0D, 0x1E, 0x35), border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.8, 6.1, 9.7, 0.45,
             "CI/CD automatique  •  SSL partout  •  99,9 % uptime  •  Backups quotidiens",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_footer(s, 16, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 17: RÉSULTATS & MÉTRIQUES
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "11")
add_slide_title(s, "Résultats et Métriques Clés")

# Stat cards - 2 rows of 4
stats = [
    ("95", "LOCALISATIONS\nGÉOLocalisées", OCP_GREEN),
    ("3 423", "NŒUDS DE\nNAVIGATION", RGBColor(0x3A, 0x8A, 0xD6)),
    ("13 346", "ARÊTES\n(CHEMINS)", RGBColor(0xE8, 0x8A, 0x2E)),
    ("100 %", "CONNECTIVITÉ\nA*", OCP_GREEN),
    ("< 3 s", "TEMPS DE\nROUTAGE", RGBColor(0x3A, 0x8A, 0xD6)),
    ("14,4 px", "PRÉCISION\nCALIBRATION", OCP_GREEN),
    ("3", "LANGUES\n(EN/FR/AR)", RGBColor(0x9B, 0x59, 0xB6)),
    ("19", "DASHBOARDS\nPERSONNALISÉS", RGBColor(0xE8, 0x3E, 0x3E)),
]

card_w = 2.7
card_h = 2.1
gap_x = 0.3
gap_y = 0.35
x_start = 0.8
y_start = 2.1

for i, (num, lbl, accent) in enumerate(stats):
    col = i % 4
    row = i // 4
    cx = x_start + col * (card_w + gap_x)
    cy = y_start + row * (card_h + gap_y)
    
    add_stat_card(s, cx, cy, card_w, card_h, num, lbl, accent=accent)

add_footer(s, 17, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 18: DÉFIS TECHNIQUES
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "12")
add_slide_title(s, "Défis Techniques et Solutions")

challenges = [
    ("Défi", "Solution"),
    ("GPS imprécis en milieu industriel\n(murs, obstructions)", "Transformation affine\navec 8 points de contrôle\nRMS = 14,4 px"),
    ("Performance Canvas avec 95+\nmarqueurs et labels", "Rendu DPR-aware,\nlabels adaptatifs,\nvirtualisation des labels"),
    ("A* sur 3 423 nœuds\net 13 346 arêtes", "Tas binaire,\nadjacence paresseuse,\noptimisation mémoire"),
    ("Multi-langue avec RTL arabe\net 430+ clés i18n", "Provider React custom,\n430+ clés,\nadaptation directionnelle"),
    ("Sécurité multi-niveaux\n(19 profils, 76 comptes)", "JWT double token,\nArgon2, middleware authorize,\nZod + rate limiting"),
]

row_h = 0.85
x_start = 0.8
y_start = 2.0
col_widths = [5.6, 6.1]
col_x = [x_start, x_start + col_widths[0] + 0.5]

for row_idx, row_data in enumerate(challenges):
    y = y_start + row_idx * row_h
    is_header = row_idx == 0
    bg_color = RGBColor(0x0D, 0x1E, 0x35) if is_header else (CARD_BG if row_idx % 2 == 1 else RGBColor(0x10, 0x22, 0x3C))
    
    for col_idx, text in enumerate(row_data):
        add_rect(s, col_x[col_idx], y, col_widths[col_idx], row_h, bg_color,
                 border_color=RGBColor(0x1E, 0x3A, 0x5C))
        font_color = OCP_GREEN if is_header else (AMBER if col_idx == 0 else OCP_GREEN)
        add_text_box(s, col_x[col_idx] + 0.15, y + 0.05, col_widths[col_idx] - 0.3, row_h - 0.1,
                     text, font_size=10 if not is_header else 10.5, bold=is_header, color=font_color)

add_footer(s, 18, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 19: LIMITES & PERSPECTIVES
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
add_section_number(s, "13")
add_slide_title(s, "Limites et Perspectives")

# Left: Limites
add_rounded_rect(s, 0.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 1.1, 2.3, 5.0, 0.4, "Limites Actuelles",
             font_size=16, bold=True, color=AMBER, font_name="Calibri")
add_rect(s, 1.1, 2.75, 1.5, 0.04, AMBER)

limits = [
    "Précision GPS limitée en milieu industriel\n(14,4 px ≈ 100 m à l'échelle actuelle)",
    "Calibration nécessaire si l'image\nde la carte change",
    "Pas d'intégration avec les systèmes\nde gestion internes OCP",
    "Mode hors-ligne non disponible",
    "Pas de gestion des événements\net maintenance en temps réel",
]
for i, txt in enumerate(limits):
    add_text_box(s, 1.3, 3.0 + i * 0.72, 4.8, 0.7, f"▸  {txt}",
                 font_size=11, color=LIGHT_GRAY, line_spacing=15)

# Right: Perspectives
add_rounded_rect(s, 6.8, 2.1, 5.6, 4.6, CARD_BG, border_color=RGBColor(0x1E, 0x3A, 0x5C))
add_text_box(s, 7.1, 2.3, 5.0, 0.4, "Évolutions Futures",
             font_size=16, bold=True, color=OCP_GREEN, font_name="Calibri")
add_green_accent_line(s, 7.1, 2.75, 1.5)

perspectives = [
    "Intégration API avec les systèmes\nde gestion internes OCP",
    "Application mobile native\n(iOS / Android) avec AR",
    "Mode hors-ligne avec cache\net synchronisation",
    "Analytiques et statistiques\nd'utilisation avancées",
    "Intégration IoT pour capteurs\nde présence en temps réel",
]
for i, txt in enumerate(perspectives):
    add_text_box(s, 7.3, 3.0 + i * 0.72, 4.8, 0.7, f"▸  {txt}",
                 font_size=11, color=LIGHT_GRAY, line_spacing=15)

add_footer(s, 19, TOTAL_SLIDES)


# ────────────────────────────────────────────────────────────────────────
# SLIDE 20: CONCLUSION
# ────────────────────────────────────────────────────────────────────────
s = add_slide(prs)
# Full dark background
add_rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, NAVY)
# Green accent line
add_rect(s, 0, 3.4, SLIDE_W_IN, 0.06, OCP_GREEN)

# Title
add_text_box(s, 1.0, 0.6, 11, 0.7, "Conclusion",
             font_size=34, bold=True, color=WHITE, font_name="Calibri Light",
             alignment=PP_ALIGN.CENTER)

# Key achievements - two columns
col1 = [
    "Projet fonctionnel et déployé",
    "95 localisations géolocalisées",
    "Carte interactive Canvas HTML5",
    "Recherche et routage A*",
]
col2 = [
    "19 dashboards personnalisés",
    "Sécurité JWT + Argon2 + Zod",
    "Trilingue EN/FR/AR (RTL)",
    "Application web complète",
]
for i, txt in enumerate(col1):
    add_text_box(s, 1.5, 1.6 + i * 0.5, 5, 0.45, f"✓  {txt}",
                 font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)
for i, txt in enumerate(col2):
    add_text_box(s, 7.0, 1.6 + i * 0.5, 5, 0.45, f"✓  {txt}",
                 font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

# Green accent line
add_rect(s, 3.5, 3.8, 6.3, 0.04, OCP_GREEN)

# Thank you section
add_text_box(s, 1.0, 4.1, 11, 0.6, "Merci pour votre attention",
             font_size=26, bold=True, color=WHITE, font_name="Calibri Light",
             alignment=PP_ALIGN.CENTER)

# Questions
add_rounded_rect(s, 5.0, 4.9, 3.3, 0.65, OCP_GREEN)
add_text_box(s, 5.0, 4.96, 3.3, 0.52, "Questions ?",
             font_size=19, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
             font_name="Calibri")

# Footer
add_rect(s, 0, 7.08, SLIDE_W_IN, 0.42, RGBColor(0x06, 0x0E, 0x1A))
add_text_box(s, 0.5, 7.14, 6, 0.3, "OCP eGuide  •  2026",
             font_size=9, color=MID_GRAY)
add_text_box(s, 11.5, 7.14, 1.5, 0.3, f"20 / {TOTAL_SLIDES}",
             font_size=9, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════
output_path = "OCP_eGuide_Soutenance_v4.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Dimensions: 16:9 widescreen (13.33\" x 7.5\")")
