#!/usr/bin/env python3
"""
ALL-IN-ONE OCP eGuide PowerPoint Polisher — v2
Uses proper add_picture API for logos.
"""
import sys
sys.path.insert(0, '/home/imrane/.pyenv/versions/3.12.11/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Emu, Inches, Cm
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
import io, os

INPUT  = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3.pptx"
OUTPUT = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3_POLISHED.pptx"
LOGO   = "/mnt/c/Users/unknown/Downloads/project assets/ocp_logo_for_doc.png"

SLIDE_W_EMU = 12191695
LOGO_CM = 1.2  # 1.2cm square

print("=" * 60)
print("OCP eGuide PowerPoint Polisher — v2")
print("=" * 60)

prs = Presentation(INPUT)
print(f"Loaded: {len(prs.slides)} slides")

# ============================================================
# 1. ADD OCP LOGO TO EVERY SLIDE (top-right, small)
# ============================================================
print("\n[1/6] Adding OCP logo...")
logo_count = 0
logo_stream = io.BytesIO(open(LOGO, 'rb').read())

# Calculate position: right-aligned with 0.75cm margin, 0.5cm from top
logo_left_cm = (SLIDE_W_EMU / 360000) - LOGO_CM - 0.75  # right edge minus logo minus margin
logo_top_cm = 0.5

for si, slide in enumerate(prs.slides):
    try:
        # Use add_picture with in-memory stream
        slide.shapes.add_picture(
            logo_stream,
            Cm(logo_left_cm),
            Cm(logo_top_cm),
            Cm(LOGO_CM),
            Cm(LOGO_CM),
        )
        logo_count += 1
        # Reset stream for next slide
        logo_stream.seek(0)
    except Exception as ex:
        print(f"  ⚠️ Slide {si+1}: {ex}")
        logo_stream.seek(0)

print(f"  ✅ Logo added to {logo_count}/{len(prs.slides)} slides")

# ============================================================
# 2. REPLACE FOOTER TEXT
# ============================================================
print("\n[2/6] Fixing footer text...")
footer_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Rapport de Stage 2026" in run.text:
                    run.text = run.text.replace("Rapport de Stage 2026", "Présentation 2026")
                    footer_count += 1
print(f"  ✅ Fixed {footer_count} footer instances")

# ============================================================
# 3. FIX SLIDE 1 TITLE
# ============================================================
print("\n[3/6] Fixing slide 1 title...")
for shape in prs.slides[0].shapes:
    if not hasattr(shape, "text_frame"):
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if "Rapport de Stage — Projet de Fin d'Études" in run.text:
                run.text = run.text.replace(
                    "Rapport de Stage — Projet de Fin d'Études",
                    "Présentation — Projet de Fin d'Études"
                )
                print("  ✅ Fixed")

# ============================================================
# 4. REMOVE ENCADRANT NAME
# ============================================================
print("\n[4/6] Removing encadrant name...")
for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Karim Alaoui" in run.text:
                    run.text = run.text.replace("M. Karim Alaoui", "______________________")
                    print(f"  ✅ Blanked encadrant in '{shape.name}'")
                elif "Karim" in run.text:
                    run.text = run.text.replace("Karim Alaoui", "")
                    run.text = run.text.replace("M. Karim", "______________________")

# ============================================================
# 5. FIX BROKEN IMAGE ON SLIDE 13
# ============================================================
print("\n[5/6] Fixing broken image...")
slide13 = prs.slides[12]
for shape in list(slide13.shapes):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            _ = shape.image
            print(f"  ✅ '{shape.name}' is OK")
        except Exception:
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"  ✅ Removed broken '{shape.name}'")

# ============================================================
# 6. SET LANGUAGE PROOFING
# ============================================================
print("\n[6/6] Setting language proofing...")
lang_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        rPr.set(qn('a:lang'), 'fr-FR')
                        lang_count += 1
        if hasattr(shape, 'has_table') and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            rPr = run._r.find(qn('a:rPr'))
                            if rPr is not None:
                                rPr.set(qn('a:lang'), 'fr-FR')
                                lang_count += 1
print(f"  ✅ Set language on {lang_count} text runs")

# ============================================================
# VERIFICATION
# ============================================================
print("\n" + "=" * 60)
print("VERIFICATION")
print("=" * 60)

rapport_ok = True
enc_ok = True
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            txt = shape.text_frame.text
            if "Rapport" in txt:
                print(f"  ⚠️ 'Rapport' still in '{shape.name}'")
                rapport_ok = False
            if "Karim" in txt or "Alaoui" in txt:
                print(f"  ⚠️ Encadrant name still in '{shape.name}'")
                enc_ok = False

if rapport_ok:
    print("  ✅ No 'Rapport' references remain")
if enc_ok:
    print("  ✅ No encadrant name references remain")

pres_count = sum(1 for s in prs.slides for sh in s.shapes
                 if hasattr(sh, "text_frame") and "Présentation 2026" in sh.text_frame.text)
print(f"  ✅ 'Présentation 2026' found {pres_count} times")

img_count = sum(1 for s in prs.slides for sh in s.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
print(f"  ✅ Total images: {img_count}")

# ============================================================
# SAVE
# ============================================================
print("\n" + "=" * 60)
print("Saving...")
prs.save(OUTPUT)
print(f"✅ Saved to: {OUTPUT}")
print(f"File size: {os.path.getsize(OUTPUT)/1024/1024:.1f} MB")
print("\n✅ ALL DONE!")
