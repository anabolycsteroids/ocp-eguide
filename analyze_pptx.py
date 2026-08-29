#!/usr/bin/env python3
"""Deep analysis of the PowerPoint presentation."""
import sys
sys.path.insert(0, '/home/imrane/.pyenv/versions/3.12.11/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
import json

PPTX_PATH = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3.pptx"

prs = Presentation(PPTX_PATH)
slide_width = prs.slide_width
slide_height = prs.slide_height

print("=" * 70)
print("POWERPOINT ANALYSIS")
print("=" * 70)
print(f"Slide dimensions: {slide_width} EMU x {slide_height} EMU")
print(f"  = {slide_width / 914400:.2f} inches x {slide_height / 914400:.2f} inches")
print(f"  = {slide_width / 360000:.1f} cm x {slide_height / 360000:.1f} cm")
print(f"Total slides: {len(prs.slides)}")

# Analyze each slide
for slide_idx, slide in enumerate(prs.slides):
    print(f"\n{'='*70}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'='*70}")
    
    # Get slide layout name
    if slide.slide_layout:
        print(f"  Layout: {slide.slide_layout.name}")
    
    shapes = list(slide.shapes)
    print(f"  Shapes: {len(shapes)}")
    
    for s_idx, shape in enumerate(shapes):
        shape_type = shape.shape_type
        name = shape.name
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        # Convert to cm for readability
        left_cm = left / 360000
        top_cm = top / 360000
        width_cm = width / 360000
        height_cm = height / 360000
        
        print(f"\n  Shape {s_idx}: '{name}'")
        print(f"    Type: {shape_type}")
        print(f"    Position: left={left_cm:.1f}cm, top={top_cm:.1f}cm")
        print(f"    Size: {width_cm:.1f}cm x {height_cm:.1f}cm")
        
        # Check if shape extends beyond slide
        if left_cm + width_cm > slide_width / 360000 + 1:
            print(f"    ⚠️ EXTENDS BEYOND RIGHT EDGE!")
        if top_cm + height_cm > slide_height / 360000 + 1:
            print(f"    ⚠️ EXTENDS BEYOND BOTTOM EDGE!")
        if left_cm < -1:
            print(f"    ⚠️ EXTENDS BEYOND LEFT EDGE!")
        if top_cm < -1:
            print(f"    ⚠️ EXTENDS BEYOND TOP EDGE!")
        
        if hasattr(shape, "text") and shape.text:
            text_preview = shape.text.strip()[:120]
            print(f"    Text: '{text_preview}'")
        
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            print(f"    Image: content_type={image.content_type}, size={len(image.blob)} bytes")
        
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"    Group with {len(shape.shapes)} sub-shapes")
            for sub in shape.shapes:
                sub_left = (left + sub.left) / 360000
                sub_top = (top + sub.top) / 360000
                sub_w = sub.width / 360000
                sub_h = sub.height / 360000
                sub_text = ""
                if hasattr(sub, "text") and sub.text:
                    sub_text = f" text='{sub.text.strip()[:60]}'"
                print(f"      Sub: type={sub.shape_type} pos=({sub_left:.1f},{sub_top:.1f}) size=({sub_w:.1f}x{sub_h:.1f}){sub_text}")

    # Check for "Rapport" or "Karim" mentions
    for shape in shapes:
        if hasattr(shape, "text") and shape.text:
            if "Rapport" in shape.text:
                print(f"\n  ⚠️ FOUND 'Rapport' in shape '{shape.name}': '{shape.text.strip()[:100]}'")
            if "Karim" in shape.text or "Alaoui" in shape.text:
                print(f"\n  ⚠️ FOUND ENCADRANT NAME in shape '{shape.name}': '{shape.text.strip()[:100]}'")

print(f"\n{'='*70}")
print("SUMMARY")
print(f"{'='*70}")
print(f"Total slides: {len(prs.slides)}")

# Count images
total_images = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            total_images += 1
print(f"Total images: {total_images}")

# Check for Rapport references
print("\nSearching for 'Rapport' references...")
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text") and "Rapport" in shape.text:
            print(f"  Slide {slide_idx+1}: '{shape.text.strip()[:100]}'")

print("\nSearching for encadrant name...")
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text") and ("Karim" in shape.text or "Alaoui" in shape.text):
            print(f"  Slide {slide_idx+1}: '{shape.text.strip()[:100]}'")
