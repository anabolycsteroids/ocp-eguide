#!/usr/bin/env python3
"""
ALL-IN-ONE OCP eGuide PowerPoint Polisher
Works from the ORIGINAL v3 file. Applies ALL fixes in one pass.
"""
import sys
sys.path.insert(0, '/home/imrane/.pyenv/versions/3.12.11/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.opc.package import Part, PackURI
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
import os

# Paths
INPUT  = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3.pptx"
OUTPUT = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3_POLISHED.pptx"
LOGO   = "/mnt/c/Users/unknown/Downloads/project assets/ocp_logo_for_doc.png"

SLIDE_W = 12191695
LOGO_SIZE = int(Emu(432000))  # ~1.2cm square
LOGO_LEFT = int(SLIDE_W - LOGO_SIZE - Emu(270000))
LOGO_TOP  = int(Emu(180000))

print("=" * 60)
print("OCP eGuide PowerPoint Polisher — ALL-IN-ONE")
print("=" * 60)

prs = Presentation(INPUT)
print(f"Loaded: {len(prs.slides)} slides")

with open(LOGO, 'rb') as f:
    logo_data = f.read()

# ============================================================
# 1. ADD OCP LOGO TO EVERY SLIDE
# ============================================================
print("\n[1/6] Adding OCP logo...")
logo_count = 0
for si, slide in enumerate(prs.slides):
    try:
        img_name = PackURI(f'/ppt/media/ocp_logo_s{si+1}.png')
        img_part = Part(img_name, 'image/png', logo_data, prs.part.package)
        rid = slide.part.relate_to(img_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')

        sp = etree.SubElement(slide.shapes._spTree, qn('p:pic'))
        nvPicPr = etree.SubElement(sp, qn('p:nvPicPr'))
        cNvPr = etree.SubElement(nvPicPr, qn('p:cNvPr'))
        cNvPr.set('id', str(9000 + si))
        cNvPr.set('name', 'OCP_Logo')
        cNvPr.set('descr', 'OCP Logo')
        etree.SubElement(nvPicPr, qn('p:cNvPicPr')).append(etree.SubElement(etree.Element('x'), qn('a:picLocks'))).getparent().find(qn('a:picLocks')).set('noChangeAspect', '1')
        etree.SubElement(nvPicPr, qn('p:nvPr'))

        bf = etree.SubElement(sp, qn('p:blipFill'))
        bf.set('rotWithShape', '0')
        blip = etree.SubElement(bf, qn('a:blip'))
        blip.set(qn('r:embed'), rid)
        st = etree.SubElement(bf, qn('a:stretch'))
        etree.SubElement(st, qn('a:fillRect'))

        spPr = etree.SubElement(sp, qn('p:spPr'))
        xf = etree.SubElement(spPr, qn('a:xfrm'))
        o = etree.SubElement(xf, qn('a:off')); o.set('x', str(LOGO_LEFT)); o.set('y', str(LOGO_TOP))
        e = etree.SubElement(xf, qn('a:ext')); e.set('cx', str(LOGO_SIZE)); e.set('cy', str(LOGO_SIZE))
        pg = etree.SubElement(spPr, qn('a:prstGeom')); pg.set('prst', 'rect')
        etree.SubElement(spPr, qn('a:noFill'))
        ln = etree.SubElement(spPr, qn('a:ln'))
        etree.SubElement(ln, qn('a:noFill'))

        slide.shapes._spTree.remove(sp)
        slide.shapes._spTree.append(sp)
        logo_count += 1
    except Exception as ex:
        print(f"  ⚠️ Slide {si+1}: {ex}")

print(f"  ✅ Logo added to {logo_count}/{len(prs.slides)} slides")

# ============================================================
# 2. REPLACE "Rapport de Stage 2026" → "Présentation 2026"
# ============================================================
print("\n[2/6] Fixing footer text...")
footer_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Rapport de Stage 2026" in run.text:
                    run.text = run.text.replace("Rapport de Stage 2026", "Présentation 2026")
                    footer_count += 1
print(f"  ✅ Fixed {footer_count} footer instances")

# ============================================================
# 3. FIX SLIDE 1 TITLE
# ============================================================
print("\n[3/6] Fixing slide 1 title...")
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if not hasattr(shape, "text_frame"):
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if "Rapport de Stage — Projet de Fin d'Études" in run.text:
                run.text = run.text.replace(
                    "Rapport de Stage — Projet de Fin d'Études",
                    "Présentation — Projet de Fin d'Études"
                )
                print("  ✅ Slide 1 title fixed")

# ============================================================
# 4. REMOVE ENCADRANT NAME
# ============================================================
print("\n[4/6] Removing encadrant name...")
for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Karim Alaoui" in run.text:
                    run.text = run.text.replace("M. Karim Alaoui", "")
                    run.text = run.text.replace("Karim Alaoui", "")
                    print(f"  ✅ Removed 'Karim Alaoui' from '{shape.name}'")
                # Blank the encadrant field
                if "Encadrant" in run.text and "____" not in run.text:
                    # Check if there's a name after Encadrant :
                    parts = run.text.split("Encadrant")
                    if len(parts) > 1:
                        after = parts[1].strip()
                        if after.startswith(":"):
                            after_content = after[1:].strip()
                            if after_content and "____" not in after_content:
                                run.text = parts[0] + "Encadrant : ____________________"
                                print(f"  ✅ Blanked encadrant field")

# ============================================================
# 5. FIX BROKEN IMAGE ON SLIDE 13
# ============================================================
print("\n[5/6] Fixing broken image...")
slide13 = prs.slides[12]
for shape in list(slide13.shapes):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            _ = shape.image
            print(f"  ✅ '{shape.name}' is OK")
        except Exception:
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"  ✅ Removed broken '{shape.name}'")

# ============================================================
# 6. SET LANGUAGE PROOFING
# ============================================================
print("\n[6/6] Setting language proofing...")
lang_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        rPr.set(qn('a:lang'), 'fr-FR')
                        lang_count += 1
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            rPr = run._r.find(qn('a:rPr'))
                            if rPr is not None:
                                rPr.set(qn('a:lang'), 'fr-FR')
                                lang_count += 1
print(f"  ✅ Set language on {lang_count} text runs")

# ============================================================
# FINAL VERIFICATION
# ============================================================
print("\n" + "=" * 60)
print("VERIFICATION")
print("=" * 60)

# Check Rapport
rapport_found = False
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            if "Rapport" in shape.text_frame.text:
                print(f"  ⚠️ 'Rapport' still in '{shape.name}': {shape.text_frame.text.strip()[:60]}")
                rapport_found = True
if not rapport_found:
    print("  ✅ No 'Rapport' references remain")

# Check encadrant
enc_found = False
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            if "Karim" in shape.text_frame.text or "Alaoui" in shape.text_frame.text:
                print(f"  ⚠️ Encadrant name still in '{shape.name}'")
                enc_found = True
if not enc_found:
    print("  ✅ No encadrant name references remain")

# Check Présentation 2026
pres_count = sum(1 for s in prs.slides for sh in s.shapes
                 if hasattr(sh, "text_frame") and "Présentation 2026" in sh.text_frame.text)
print(f"  ✅ 'Présentation 2026' found {pres_count} times")

# Count images
img_count = sum(1 for s in prs.slides for sh in s.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
print(f"  ✅ Total images: {img_count} (original + logos)")

# ============================================================
# SAVE
# ============================================================
print("\n" + "=" * 60)
print("Saving...")
prs.save(OUTPUT)
print(f"✅ Saved to: {OUTPUT}")
print(f"File size: {os.path.getsize(OUTPUT)/1024/1024:.1f} MB")
print("\n✅ ALL DONE!")
