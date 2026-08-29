#!/usr/bin/env python3
"""Add OCP logo to every slide - per-slide image approach."""
import sys
sys.path.insert(0, '/home/imrane/.pyenv/versions/3.12.11/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.opc.package import Part, PackURI
from lxml import etree
import os

PPTX_PATH = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3_POLISHED.pptx"
PPTX_OUTPUT = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3_POLISHED.pptx"
LOGO_PATH = "/mnt/c/Users/unknown/Downloads/project assets/ocp_logo_for_doc.png"

SLIDE_W = 12191695
LOGO_SIZE = Emu(432000)  # ~1.2cm
LOGO_LEFT = int(SLIDE_W - LOGO_SIZE - Emu(270000))
LOGO_TOP = int(Emu(180000))

print("Loading presentation...")
prs = Presentation(PPTX_PATH)

with open(LOGO_PATH, 'rb') as f:
    logo_data = f.read()

logo_count = 0
for slide_idx, slide in enumerate(prs.slides):
    try:
        # Add image part to this slide
        image_part_name = PackURI(f'/ppt/media/ocp_logo_s{slide_idx+1}.png')
        image_part = Part(image_part_name, 'image/png', logo_data, prs.part.package)
        
        slide_rId = slide.part.relate_to(
            image_part,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'
        )
        
        # Create pic element
        sp = etree.SubElement(slide.shapes._spTree, qn('p:pic'))
        
        nvPicPr = etree.SubElement(sp, qn('p:nvPicPr'))
        cNvPr = etree.SubElement(nvPicPr, qn('p:cNvPr'))
        cNvPr.set('id', str(9000 + slide_idx))
        cNvPr.set('name', 'OCP_Logo')
        cNvPr.set('descr', 'OCP Logo')
        
        cNvPicPr = etree.SubElement(nvPicPr, qn('p:cNvPicPr'))
        picLocks = etree.SubElement(cNvPicPr, qn('a:picLocks'))
        picLocks.set('noChangeAspect', '1')
        
        nvPr = etree.SubElement(nvPicPr, qn('p:nvPr'))
        
        blipFill = etree.SubElement(sp, qn('p:blipFill'))
        blipFill.set('rotWithShape', '0')
        blip = etree.SubElement(blipFill, qn('a:blip'))
        blip.set(qn('r:embed'), slide_rId)
        stretch = etree.SubElement(blipFill, qn('a:stretch'))
        etree.SubElement(stretch, qn('a:fillRect'))
        
        spPr = etree.SubElement(sp, qn('p:spPr'))
        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
        off = etree.SubElement(xfrm, qn('a:off'))
        off.set('x', str(LOGO_LEFT))
        off.set('y', str(LOGO_TOP))
        ext = etree.SubElement(xfrm, qn('a:ext'))
        ext.set('cx', str(int(LOGO_SIZE)))
        ext.set('cy', str(int(LOGO_SIZE)))
        prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        prstGeom.set('prst', 'rect')
        etree.SubElement(spPr, qn('a:noFill'))
        ln = etree.SubElement(spPr, qn('a:ln'))
        etree.SubElement(ln, qn('a:noFill'))
        
        # Move to front
        slide.shapes._spTree.remove(sp)
        slide.shapes._spTree.append(sp)
        
        logo_count += 1
    except Exception as e:
        print(f"  ⚠️ Slide {slide_idx+1}: {e}")

print(f"Added logo to {logo_count}/{len(prs.slides)} slides")

# Save to a temp file first to avoid overwriting
temp_output = PPTX_OUTPUT + ".tmp"
prs.save(temp_output)
os.replace(temp_output, PPTX_OUTPUT)

print(f"✅ Saved to: {PPTX_OUTPUT}")
print(f"File size: {os.path.getsize(PPTX_OUTPUT)/1024/1024:.1f} MB")
