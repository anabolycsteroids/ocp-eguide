#!/usr/bin/env python3
"""Add OCP logo to every slide using direct XML manipulation."""
import sys
sys.path.insert(0, '/home/imrane/.pyenv/versions/3.12.11/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.opc.package import Part, PackURI
from lxml import etree
import os

PPTX_PATH = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3_POLISHED.pptx"
PPTX_OUTPUT = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3_POLISHED.pptx"
LOGO_PATH = "/mnt/c/Users/unknown/Downloads/project assets/ocp_logo_for_doc.png"

# Slide dimensions
SLIDE_W = 12191695
SLIDE_H = 6858000

# Logo: small, subtle, top-right
LOGO_SIZE = Emu(432000)  # ~1.2cm
LOGO_LEFT = Emu(SLIDE_W - LOGO_SIZE - Emu(270000))  # ~0.75cm from right
LOGO_TOP = Emu(180000)  # ~0.5cm from top

print("Loading presentation...")
prs = Presentation(PPTX_PATH)

# Read logo file
with open(LOGO_PATH, 'rb') as f:
    logo_data = f.read()

print(f"Logo data: {len(logo_data)} bytes")

# Create a shared image part that all slides will reference
image_part_name = PackURI('/ppt/media/ocp_logo.png')
image_part = Part(
    image_part_name,
    'image/png',
    logo_data,
    prs.part.package
)

# Add relationship from presentation part to image
rId = prs.part.relate_to(
    image_part,
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'
)
print(f"Image relationship ID: {rId}")

# Now add logo to each slide
logo_count = 0
for slide_idx, slide in enumerate(prs.slides):
    try:
        # We need to add a relationship from the slide to the same image part
        # Each slide needs its own relationship to the shared image part
        slide_rId = slide.part.relate_to(
            image_part,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'
        )
        
        # Create the pic element via XML
        sp = etree.SubElement(slide.shapes._spTree, qn('p:pic'))
        
        # Non-visual properties
        nvPicPr = etree.SubElement(sp, qn('p:nvPicPr'))
        
        cNvPr = etree.SubElement(nvPicPr, qn('p:cNvPr'))
        cNvPr.set('id', str(9000 + slide_idx))
        cNvPr.set('name', f'OCP_Logo')
        cNvPr.set('descr', 'OCP Logo')
        cNvPr.set('title', 'OCP Logo')
        
        cNvPicPr = etree.SubElement(nvPicPr, qn('p:cNvPicPr'))
        picLocks = etree.SubElement(cNvPicPr, qn('a:picLocks'))
        picLocks.set('noChangeAspect', '1')
        
        nvPr = etree.SubElement(nvPicPr, qn('p:nvPr'))
        
        # BlipFill
        blipFill = etree.SubElement(sp, qn('p:blipFill'))
        blipFill.set('rotWithShape', '0')
        
        blip = etree.SubElement(blipFill, qn('a:blip'))
        blip.set(qn('r:embed'), slide_rId)
        
        # Soft edge effect for subtle look
        softEdge = etree.SubElement(blip, qn('a:softEdge'))
        softEdge.set('rad', '12700')  # 0.1cm soft edge
        
        stretch = etree.SubElement(blipFill, qn('a:stretch'))
        fillRect = etree.SubElement(stretch, qn('a:fillRect'))
        
        # Shape properties
        spPr = etree.SubElement(sp, qn('p:spPr'))
        
        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
        off = etree.SubElement(xfrm, qn('a:off'))
        off.set('x', str(int(LOGO_LEFT)))
        off.set('y', str(int(LOGO_TOP)))
        ext = etree.SubElement(xfrm, qn('a:ext'))
        ext.set('cx', str(int(LOGO_SIZE)))
        ext.set('cy', str(int(LOGO_SIZE)))
        
        prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        prstGeom.set('prst', 'rect')
        
        # No fill
        etree.SubElement(spPr, qn('a:noFill'))
        
        # No outline
        ln = etree.SubElement(spPr, qn('a:ln'))
        etree.SubElement(ln, qn('a:noFill'))
        
        # Move to front (top of z-order)
        slide.shapes._spTree.remove(sp)
        slide.shapes._spTree.append(sp)
        
        logo_count += 1
        
    except Exception as e:
        print(f"  ⚠️ Slide {slide_idx+1}: {e}")
        import traceback
        traceback.print_exc()

print(f"Added logo to {logo_count}/{len(prs.slides)} slides")

# Save
print("Saving...")
prs.save(PPTX_OUTPUT)
print(f"✅ Saved to: {PPTX_OUTPUT}")
print(f"File size: {os.path.getsize(PPTX_OUTPUT)/1024/1024:.1f} MB")
