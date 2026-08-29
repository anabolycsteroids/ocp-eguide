#!/usr/bin/env python3
"""
OCP eGuide PowerPoint Polisher
Preserves ALL existing content. Only fixes:
1. Adds small OCP logo to each slide
2. Replaces "Rapport de Stage 2026" → "Présentation 2026"
3. Removes encadrant name
4. Fixes broken image on slide 13
5. Checks image positioning
"""
import sys
sys.path.insert(0, '/home/imrane/.pyenv/versions/3.12.11/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import copy
import os

# Paths
PPTX_INPUT = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3.pptx"
PPTX_OUTPUT = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3_POLISHED.pptx"
LOGO_PATH = "/mnt/c/Users/unknown/Downloads/project assets/ocp_logo_for_doc.png"

# Slide dimensions
SLIDE_W = 12191695  # EMU (33.867 cm = 13.33 inches)
SLIDE_H = 6858000   # EMU (19.05 cm = 7.5 inches)

# Logo size (small, subtle): ~1.2cm tall
LOGO_H = Emu(432000)  # 1.2cm
LOGO_W = Emu(432000)  # Square-ish for logo
# Logo position: top-right corner with small margin
LOGO_LEFT = Emu(SLIDE_W - LOGO_W - Emu(360000))  # 1cm from right
LOGO_TOP = Emu(180000)  # 0.5cm from top

print("=" * 60)
print("OCP eGuide PowerPoint Polisher")
print("=" * 60)

prs = Presentation(PPTX_INPUT)
print(f"Loaded: {len(prs.slides)} slides")

# ============================================================
# STEP 1: Add OCP logo to every slide
# ============================================================
print("\n[1/8] Adding OCP logo to every slide...")

# First, we need to add the image part to the presentation
from pptx.opc.package import Part, PackURI

# Read logo file
with open(LOGO_PATH, 'rb') as f:
    logo_data = f.read()

# We'll add the image as a relationship to the first slide, then reference it
# Actually, python-pptx handles this per-slide. Let's add to each slide.

logo_added_count = 0
for slide_idx, slide in enumerate(prs.slides):
    try:
        # Add picture to slide
        # We need to add the image part and relationship
        
        # Use the slide's part to add the image
        image_part_name = PackURI(f'/ppt/media/ocp_logo_s{slide_idx+1}.png')
        
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.image.image import Image as PptxImage
        
        logo_part = Part(
            image_part_name,
            'image/png',
            logo_data,
            prs.part.package
        )
        
        # Add relationship from slide to image part
        rId = slide.part.relate_to(logo_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')
        
        # Create the picture shape via XML manipulation
        sp = etree.SubElement(slide.shapes._spTree, qn('p:pic'))
        
        # Non-visual picture properties
        nvPicPr = etree.SubElement(sp, qn('p:nvPicPr'))
        cNvPr = etree.SubElement(nvPicPr, qn('p:cNvPr'))
        cNvPr.set('id', str(1000 + slide_idx))
        cNvPr.set('name', f'OCP_Logo_{slide_idx+1}')
        cNvPr.set('descr', 'OCP Logo')
        
        cNvPicPr = etree.SubElement(nvPicPr, qn('p:cNvPicPr'))
        picLocks = etree.SubElement(cNvPicPr, qn('a:picLocks'))
        picLocks.set('noChangeAspect', '1')
        
        nvPr = etree.SubElement(nvPicPr, qn('p:nvPr'))
        
        # Picture fill
        blipFill = etree.SubElement(sp, qn('p:blipFill'))
        blipFill.set('rotWithShape', '0')
        
        blip = etree.SubElement(blipFill, qn('a:blip'))
        blip.set(qn('r:embed'), rId)
        
        stretch = etree.SubElement(blipFill, qn('a:stretch'))
        fillRect = etree.SubElement(stretch, qn('a:fillRect'))
        
        # Shape properties
        spPr = etree.SubElement(sp, qn('p:spPr'))
        
        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
        off = etree.SubElement(xfrm, qn('a:off'))
        off.set('x', str(LOGO_LEFT))
        off.set('y', str(LOGO_TOP))
        ext = etree.SubElement(xfrm, qn('a:ext'))
        ext.set('cx', str(LOGO_W))
        ext.set('cy', str(LOGO_H))
        
        prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        prstGeom.set('prst', 'rect')
        
        # Add no fill to avoid background
        noFill = etree.SubElement(spPr, qn('a:noFill'))
        
        # Add line (no outline)
        ln = etree.SubElement(spPr, qn('a:ln'))
        noFill2 = etree.SubElement(ln, qn('a:noFill'))
        
        # Set z-order to be on top (move to end of spTree)
        slide.shapes._spTree.remove(sp)
        slide.shapes._spTree.append(sp)
        
        logo_added_count += 1
    except Exception as e:
        print(f"  ⚠️ Failed to add logo to slide {slide_idx+1}: {e}")

print(f"  Added logo to {logo_added_count}/{len(prs.slides)} slides")

# ============================================================
# STEP 2: Replace "Rapport de Stage 2026" in footers
# ============================================================
print("\n[2/8] Replacing footer text...")

footer_count = 0
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                if "Rapport de Stage 2026" in para.text:
                    # Replace in each run
                    for run in para.runs:
                        if "Rapport de Stage 2026" in run.text:
                            run.text = run.text.replace("Rapport de Stage 2026", "Présentation 2026")
                            footer_count += 1
                    # Also check direct text
                    if "Rapport de Stage 2026" in para.text:
                        # Handle case where replacement didn't happen in runs
                        for run in para.runs:
                            if "Rapport" in run.text:
                                run.text = run.text.replace("Rapport de Stage 2026", "Présentation 2026")
                                footer_count += 1

print(f"  Replaced {footer_count} footer instances")

# ============================================================
# STEP 3: Fix slide 1 title
# ============================================================
print("\n[3/8] Fixing slide 1 title...")

slide1 = prs.slides[0]
for shape in slide1.shapes:
    if hasattr(shape, "text_frame"):
        for para in shape.text_frame.paragraphs:
            if "Rapport de Stage — Projet de Fin d'Études" in para.text:
                for run in para.runs:
                    if "Rapport de Stage" in run.text:
                        run.text = run.text.replace(
                            "Rapport de Stage — Projet de Fin d'Études",
                            "Présentation — Projet de Fin d'Études"
                        )
                        print(f"  Fixed slide 1 title")

# ============================================================
# STEP 4: Remove encadrant name from slide 1
# ============================================================
print("\n[4/8] Removing encadrant name...")

for shape in slide1.shapes:
    if hasattr(shape, "text_frame"):
        full_text = shape.text_frame.text
        if "Karim" in full_text or "Alaoui" in full_text:
            print(f"  Found encadrant in shape '{shape.name}'")
            for para in shape.text_frame.paragraphs:
                if "Karim" in para.text or "Alaoui" in para.text:
                    for run in para.runs:
                        if "Karim" in run.text or "Alaoui" in run.text:
                            # Remove the name, replace with blank
                            run.text = run.text.replace("M. Karim Alaoui", "")
                            run.text = run.text.replace("Karim Alaoui", "")
                            run.text = run.text.replace("M.Karim Alaoui", "")
                            print(f"    → Removed name from run")
                    # Also check for "Encadrant" line and blank it
                    if "Encadrant" in para.text:
                        for run in para.runs:
                            if "Encadrant" in run.text:
                                # Keep "Encadrant :" but remove any name after
                                parts = run.text.split(":")
                                if len(parts) > 1:
                                    run.text = parts[0] + " : ____________________"
                                    print(f"    → Blanked encadrant field")
                # Check if the paragraph contains encadrant info
                text = para.text.strip()
                if "Encadrant" in text and "Karim" not in text and "____" not in text:
                    # Check if there's a name without "Karim" (in case name was different)
                    pass

# Also search ALL slides for encadrant name
for slide_idx, slide in enumerate(prs.slides):
    if slide_idx == 0:
        continue  # Already handled
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            if "Karim" in shape.text_frame.text or "Alaoui" in shape.text_frame.text:
                print(f"  ⚠️ Found encadrant name on slide {slide_idx+1}!")
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Karim" in run.text or "Alaoui" in run.text:
                            run.text = run.text.replace("M. Karim Alaoui", "")
                            run.text = run.text.replace("Karim Alaoui", "")

# ============================================================
# STEP 5: Check and fix image positioning
# ============================================================
print("\n[5/8] Checking image positioning...")

slide_w_cm = SLIDE_W / 360000
slide_h_cm = SLIDE_H / 360000

for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left_cm = shape.left / 360000
            top_cm = shape.top / 360000
            w_cm = shape.width / 360000
            h_cm = shape.height / 360000
            
            right_edge = left_cm + w_cm
            bottom_edge = top_cm + h_cm
            
            issues = []
            if right_edge > slide_w_cm + 0.5:
                issues.append(f"extends past right edge ({slide_w_cm:.1f}cm)")
            if bottom_edge > slide_h_cm + 0.5:
                issues.append(f"extends past bottom edge ({slide_h_cm:.1f}cm)")
            if left_cm < -0.5:
                issues.append("extends past left edge")
            if top_cm < -0.5:
                issues.append("extends past top edge")
            
            if issues:
                print(f"  ⚠️ Slide {slide_idx+1} '{shape.name}': {', '.join(issues)}")
                print(f"    Position: ({left_cm:.1f},{top_cm:.1f}) Size: ({w_cm:.1f}x{h_cm:.1f})")
            else:
                print(f"  ✅ Slide {slide_idx+1} '{shape.name}': OK ({left_cm:.1f},{top_cm:.1f}) {w_cm:.1f}x{h_cm:.1f}")

# ============================================================
# STEP 6: Fix the broken image on slide 13
# ============================================================
print("\n[6/8] Checking broken image on slide 13...")

slide13 = prs.slides[12]
for shape in slide13.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            _ = shape.image
            print(f"  ✅ '{shape.name}' is OK")
        except Exception as e:
            print(f"  ⚠️ '{shape.name}' is broken: {e}")
            print(f"    This image has a missing relationship. Removing broken shape...")
            # Remove the broken shape from the slide
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"    → Removed broken image shape")

# ============================================================
# STEP 7: Set language/proofing for spell check
# ============================================================
print("\n[7/8] Setting language proofing...")

# Set the default language to French and add English as editing language
# This is done at the presentation level via XML
try:
    # Access the presentation XML
    pres_part = prs.part
    
    # Try to set language via slide master
    # The spell check language is typically set per-run in PowerPoint
    # We'll set it on all text runs to French (1036) with English (1033) as secondary
    
    # For python-pptx, we can set the language on individual runs
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Set language to French (1036)
                        rPr = run._r.find(qn('a:rPr'))
                        if rPr is None:
                            rPr = run._r.find(qn('a:rPr'))
                        if rPr is not None:
                            rPr.set(qn('a:lang'), 'fr-FR')
            
            # Also handle tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                rPr = run._r.find(qn('a:rPr'))
                                if rPr is not None:
                                    rPr.set(qn('a:lang'), 'fr-FR')
    
    print("  Set language to fr-FR on all text runs")
except Exception as e:
    print(f"  ⚠️ Language setting failed: {e}")

# ============================================================
# STEP 8: Final verification
# ============================================================
print("\n[8/8] Final verification...")

# Check no "Rapport" remains
rapport_found = False
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            if "Rapport" in shape.text_frame.text:
                print(f"  ⚠️ 'Rapport' still found on slide {slide_idx+1}: '{shape.text_frame.text.strip()[:80]}'")
                rapport_found = True

if not rapport_found:
    print("  ✅ No 'Rapport' references remain")

# Check no encadrant name
encadrant_found = False
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            if "Karim" in shape.text_frame.text or "Alaoui" in shape.text_frame.text:
                print(f"  ⚠️ Encadrant name found on slide {slide_idx+1}")
                encadrant_found = True

if not encadrant_found:
    print("  ✅ No encadrant name references remain")

# Count shapes with "Présentation 2026"
pres_count = 0
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            if "Présentation 2026" in shape.text_frame.text:
                pres_count += 1

print(f"  ✅ 'Présentation 2026' found in {pres_count} footer locations")

# ============================================================
# SAVE
# ============================================================
print("\n" + "=" * 60)
print("Saving polished presentation...")
prs.save(PPTX_OUTPUT)
print(f"✅ Saved to: {PPTX_OUTPUT}")

# File size comparison
orig_size = os.path.getsize(PPTX_INPUT)
new_size = os.path.getsize(PPTX_OUTPUT)
print(f"Original: {orig_size/1024/1024:.1f} MB")
print(f"Polished: {new_size/1024/1024:.1f} MB")
print(f"\n✅ PowerPoint polishing complete!")
