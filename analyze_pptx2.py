#!/usr/bin/env python3
"""Continue analysis from slide 13 onwards, handling broken images."""
import sys
sys.path.insert(0, '/home/imrane/.pyenv/versions/3.12.11/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = "/mnt/c/Users/unknown/Downloads/project assets/OCP_eGuide_Defense_Final_v3.pptx"
prs = Presentation(PPTX_PATH)

# Continue from slide 13 onwards
for slide_idx in range(12, len(prs.slides)):
    slide = prs.slides[slide_idx]
    print(f"\n{'='*70}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'='*70}")
    
    shapes = list(slide.shapes)
    print(f"  Shapes: {len(shapes)}")
    
    for s_idx, shape in enumerate(shapes):
        shape_type = shape.shape_type
        name = shape.name
        left_cm = shape.left / 360000
        top_cm = shape.top / 360000
        width_cm = shape.width / 360000
        height_cm = shape.height / 360000
        slide_w = prs.slide_width / 360000
        slide_h = prs.slide_height / 360000
        
        print(f"\n  Shape {s_idx}: '{name}' type={shape_type}")
        print(f"    Pos: ({left_cm:.1f},{top_cm:.1f}) Size: ({width_cm:.1f}x{height_cm:.1f})")
        
        # Check bounds
        if left_cm + width_cm > slide_w + 0.5:
            print(f"    ⚠️ EXTENDS BEYOND RIGHT EDGE ({slide_w:.1f}cm)")
        if top_cm + height_cm > slide_h + 0.5:
            print(f"    ⚠️ EXTENDS BEYOND BOTTOM EDGE ({slide_h:.1f}cm)")
        
        if hasattr(shape, "text") and shape.text:
            print(f"    Text: '{shape.text.strip()[:120]}'")
        
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                print(f"    Image: {image.content_type}, {len(image.blob)} bytes")
            except Exception as e:
                print(f"    ⚠️ Image BROKEN: {e}")
        
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"    Group with {len(shape.shapes)} sub-shapes")
            for sub in shape.shapes:
                sub_left = (shape.left + sub.left) / 360000
                sub_top = (shape.top + sub.top) / 360000
                sub_w = sub.width / 360000
                sub_h = sub.height / 360000
                sub_text = ""
                if hasattr(sub, "text") and sub.text:
                    sub_text = f" text='{sub.text.strip()[:60]}'"
                print(f"      Sub: type={sub.shape_type} pos=({sub_left:.1f},{sub_top:.1f}) size=({sub_w:.1f}x{sub_h:.1f}){sub_text}")

# Final summary
print(f"\n{'='*70}")
print("COMPREHENSIVE SEARCH")
print(f"{'='*70}")

# Search all slides for Rapport, Karim, Alaoui
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if "Rapport" in shape.text:
                print(f"  Slide {slide_idx+1} '{shape.name}': RAPPORT → '{shape.text.strip()[:80]}'")
            if "Karim" in shape.text or "Alaoui" in shape.text:
                print(f"  Slide {slide_idx+1} '{shape.name}': ENCADRANT → '{shape.text.strip()[:80]}'")

# Count images per slide
print(f"\nImages per slide:")
for slide_idx, slide in enumerate(prs.slides):
    imgs = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if imgs:
        print(f"  Slide {slide_idx+1}: {len(imgs)} image(s)")
        for img in imgs:
            print(f"    {img.name}: ({img.left/360000:.1f},{img.top/360000:.1f}) ({img.width/360000:.1f}x{img.height/360000:.1f})")
